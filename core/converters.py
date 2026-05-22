"""基于 PyMuPDF（开源版 pymupdf）的文档转换与 PDF 加解密。"""

from __future__ import annotations

import json
import os
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

import fitz  # PyMuPDF 开源版

from progress_util import iter_progress

IMAGE_EXT = {
    ".jpg", ".jpeg", ".png", ".bmp", ".gif", ".tif", ".tiff", ".webp",
    ".jfif", ".pjpeg", ".heic", ".heif", ".ico",
}
OFFICE_EXT = {
    ".doc": "word",
    ".docx": "word",
    ".docm": "word",
    ".xls": "excel",
    ".xlsx": "excel",
    ".xlsm": "excel",
    ".ppt": "ppt",
    ".pptx": "ppt",
    ".pptm": "ppt",
}


def _open_pdf(path: Path, password: str | None = None) -> fitz.Document:
    doc = fitz.open(path)
    if doc.is_encrypted:
        pwd = password or ""
        if not doc.authenticate(pwd):
            doc.close()
            raise ValueError("PDF 已加密，请提供正确密码（-p）")
    return doc


def _default_out(src: Path, suffix: str) -> Path:
    return src.with_suffix(suffix)


def _pdf_pages(doc: fitz.Document) -> list:
    if doc.page_count <= 0:
        raise RuntimeError("PDF 没有页面或文件已损坏")
    return list(doc)


def _render_page(page: fitz.Page, dpi: int = 150) -> fitz.Pixmap:
    try:
        return page.get_pixmap(dpi=dpi, alpha=False)
    except Exception:
        zoom = dpi / 72.0
        return page.get_pixmap(matrix=fitz.Matrix(zoom, zoom), alpha=False)


def _blank_ppt_layout(prs):
    for layout in prs.slide_layouts:
        name = (layout.name or "").lower()
        if "blank" in name or "空白" in name or "empty" in name:
            return layout
    for idx in (6, 5, 1, 0):
        if idx < len(prs.slide_layouts):
            return prs.slide_layouts[idx]
    return prs.slide_layouts[0]


def pdf_to_word(src: Path, dst: Path | None = None) -> Path:
    """每页 PDF 对应 Word 中一段内容，避免重复插入内嵌图导致页数膨胀。"""
    from docx import Document
    from docx.shared import Inches

    dst = dst or _default_out(src, ".docx")
    pdf = _open_pdf(src)
    word = Document()
    tmpdir = Path(tempfile.mkdtemp(prefix="pdf2word_"))
    pages_with_content = 0
    try:
        pages = _pdf_pages(pdf)
        for i, page in enumerate(
            iter_progress(pages, desc="PDF->Word", total=len(pages), unit="页")
        ):
            if i > 0:
                word.add_page_break()

            text = page.get_text().strip()

            # 有足够文字：只写文字（不再逐张插入内嵌图，避免 6 页变 30+ 页）
            if len(text) >= 25:
                for line in text.splitlines():
                    line = line.strip()
                    if line:
                        word.add_paragraph(line)
                pages_with_content += 1
            else:
                # 图片/扫描页：只插入整页图，不加占位文字
                pix = _render_page(page, dpi=150)
                img_path = tmpdir / f"page_{i + 1:04d}.png"
                pix.save(str(img_path))
                word.add_picture(str(img_path), width=Inches(6.0))
                pages_with_content += 1

        if pages_with_content == 0:
            raise RuntimeError("PDF 无可用内容（空文档或无法解析）")

        dst.parent.mkdir(parents=True, exist_ok=True)
        word.save(dst)
    finally:
        pdf.close()
        shutil.rmtree(tmpdir, ignore_errors=True)
    return dst


def pdf_to_excel(src: Path, dst: Path | None = None) -> Path:
    import openpyxl

    dst = dst or _default_out(src, ".xlsx")
    doc = _open_pdf(src)
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Tables"
    row = 1
    found_table = False
    try:
        pages = _pdf_pages(doc)
        for page in iter_progress(pages, desc="PDF->Excel", total=len(pages), unit="页"):
            try:
                finder = page.find_tables()
            except Exception:
                finder = None
            if not finder:
                continue
            for tab in finder.tables:
                data = tab.extract()
                if not data:
                    continue
                found_table = True
                for r in data:
                    for col, cell in enumerate(r, start=1):
                        ws.cell(row=row, column=col, value=cell)
                    row += 1
                row += 1

        if not found_table:
            ws.title = "Text"
            row = 1
            for page_idx, page in enumerate(pages):
                text = page.get_text().strip()
                if text:
                    ws.cell(row=row, column=1, value=f"[第 {page_idx + 1} 页]")
                    row += 1
                    for line in text.splitlines():
                        line = line.strip()
                        if line:
                            ws.cell(row=row, column=1, value=line)
                            row += 1
                    row += 1
                else:
                    ws.cell(
                        row=row,
                        column=1,
                        value=f"[第 {page_idx + 1} 页：图片页，请用 pdf2image 导出]",
                    )
                    row += 2

        dst.parent.mkdir(parents=True, exist_ok=True)
        wb.save(dst)
    finally:
        doc.close()
    return dst


def pdf_to_ppt(src: Path, dst: Path | None = None, dpi: int = 150) -> Path:
    from pptx import Presentation
    from pptx.util import Inches

    dst = dst or _default_out(src, ".pptx")
    doc = _open_pdf(src)
    prs = Presentation()
    layout = _blank_ppt_layout(prs)
    tmpdir = Path(tempfile.mkdtemp(prefix="pdf2ppt_"))
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    try:
        pages = _pdf_pages(doc)
        for page in iter_progress(pages, desc="PDF->PPT", total=len(pages), unit="页"):
            pix = _render_page(page, dpi=dpi)
            img = tmpdir / f"p{page.number + 1:04d}.png"
            pix.save(str(img))
            slide = prs.slides.add_slide(layout)
            pic = slide.shapes.add_picture(str(img), Inches(0), Inches(0))
            pic_w, pic_h = pic.width, pic.height
            if pic_w and pic_h:
                scale = min(slide_w / pic_w, slide_h / pic_h)
                pic.width = int(pic_w * scale)
                pic.height = int(pic_h * scale)
                pic.left = int((slide_w - pic.width) / 2)
                pic.top = int((slide_h - pic.height) / 2)
        if len(prs.slides) == 0:
            raise RuntimeError("未能生成任何幻灯片")
        dst.parent.mkdir(parents=True, exist_ok=True)
        prs.save(dst)
    finally:
        doc.close()
        shutil.rmtree(tmpdir, ignore_errors=True)
    return dst


def pdf_to_images(
    src: Path,
    out_dir: Path | None = None,
    fmt: str = "png",
    dpi: int = 150,
) -> list[Path]:
    fmt = fmt.lower().lstrip(".")
    if fmt not in ("png", "jpeg", "jpg"):
        raise ValueError("图片格式仅支持 png / jpeg")
    if fmt == "jpg":
        fmt = "jpeg"

    doc = _open_pdf(src)
    out_dir = out_dir or src.parent / f"{src.stem}_pages"
    out_dir.mkdir(parents=True, exist_ok=True)
    paths: list[Path] = []
    try:
        pages = _pdf_pages(doc)
        for page in iter_progress(pages, desc="PDF->Image", total=len(pages), unit="页"):
            pix = _render_page(page, dpi=dpi)
            ext = "jpg" if fmt == "jpeg" else fmt
            out = out_dir / f"{src.stem}_p{page.number + 1:04d}.{ext}"
            pix.save(str(out))
            paths.append(out)
    finally:
        doc.close()
    if not paths:
        raise RuntimeError("未能导出任何图片")
    return paths


def pdf_to_markdown(src: Path, dst: Path | None = None) -> Path:
    dst = dst or _default_out(src, ".md")
    assets = dst.parent / f"{dst.stem}_images"
    assets.mkdir(parents=True, exist_ok=True)
    doc = _open_pdf(src)
    parts: list[str] = [f"# {src.name}\n"]
    try:
        pages = _pdf_pages(doc)
        for i, page in enumerate(
            iter_progress(pages, desc="PDF->MD", total=len(pages), unit="页")
        ):
            parts.append(f"\n## 第 {i + 1} 页\n")
            try:
                md_text = page.get_text("markdown").strip()
            except Exception:
                md_text = page.get_text().strip()
            if md_text:
                parts.append(md_text)

            img_path = assets / f"page_{i + 1:04d}.png"
            _render_page(page, dpi=120).save(str(img_path))
            rel = f"{assets.name}/{img_path.name}"
            parts.append(f"\n![第{i + 1}页]({rel})\n")
    finally:
        doc.close()
    dst.write_text("\n".join(parts), encoding="utf-8")
    print(f"Markdown 图片目录: {assets}")
    return dst


def pdf_to_json(src: Path, dst: Path | None = None) -> Path:
    dst = dst or _default_out(src, ".json")
    doc = _open_pdf(src)
    payload = {"file": str(src), "page_count": doc.page_count, "pages": []}
    try:
        pages = list(doc)
        for i, page in enumerate(
            iter_progress(pages, desc="PDF->JSON", total=len(pages), unit="页")
        ):
            payload["pages"].append(
                {
                    "index": i,
                    "width": page.rect.width,
                    "height": page.rect.height,
                    "text": page.get_text(),
                    "blocks": page.get_text("dict"),
                }
            )
    finally:
        doc.close()
    dst.write_text(
        json.dumps(payload, ensure_ascii=False, indent=2, default=str),
        encoding="utf-8",
    )
    return dst


def pdf_to_text(src: Path, dst: Path | None = None) -> Path:
    dst = dst or _default_out(src, ".txt")
    doc = _open_pdf(src)
    try:
        pages = list(doc)
        chunks = []
        for page in iter_progress(pages, desc="PDF->TXT", total=len(pages), unit="页"):
            chunks.append(page.get_text())
        text = "\n\n".join(chunks)
    finally:
        doc.close()
    dst.write_text(text, encoding="utf-8")
    return dst


def pick_image_files() -> list[Path]:
    import tkinter as tk
    from tkinter import filedialog

    root = tk.Tk()
    root.withdraw()
    root.attributes("-topmost", True)
    root.update()
    paths = filedialog.askopenfilenames(
        title="选择一张或多张图片（可多选）",
        filetypes=[
            ("图片", "*.png;*.jpg;*.jpeg;*.bmp;*.gif;*.tif;*.tiff;*.webp"),
            ("所有文件", "*.*"),
        ],
    )
    root.destroy()
    return [Path(p) for p in paths if p]


def images_to_pdf(sources: list[Path], dst: Path) -> Path:
    if not sources:
        raise ValueError("请至少选择一张图片")
    doc = fitz.open()
    try:
        for src in iter_progress(sources, desc="Image->PDF", total=len(sources), unit="张"):
            if not src.is_file():
                raise FileNotFoundError(src)
            img = fitz.open(src)
            pdf_bytes = img.convert_to_pdf()
            img.close()
            part = fitz.open("pdf", pdf_bytes)
            doc.insert_pdf(part)
            part.close()
        dst.parent.mkdir(parents=True, exist_ok=True)
        doc.save(dst)
    finally:
        doc.close()
    return dst


def _soffice_from_dir(install_dir: Path) -> Path | None:
    for sub in ("program/soffice.exe", "program/soffice.com", "soffice.exe"):
        cand = install_dir / sub.replace("/", os.sep)
        if cand.is_file():
            return cand
    return None


def _libreoffice_from_registry() -> str | None:
    if sys.platform != "win32":
        return None
    try:
        import winreg
    except ImportError:
        return None

    keys_to_try = [
        (winreg.HKEY_LOCAL_MACHINE, r"SOFTWARE\LibreOffice\Layers\LibreOffice"),
        (winreg.HKEY_CURRENT_USER, r"SOFTWARE\LibreOffice\Layers\LibreOffice"),
        (winreg.HKEY_LOCAL_MACHINE, r"SOFTWARE\LibreOffice\LibreOffice"),
    ]
    for hive, path in keys_to_try:
        try:
            with winreg.OpenKey(hive, path) as key:
                i = 0
                while True:
                    try:
                        ver = winreg.EnumKey(key, i)
                    except OSError:
                        break
                    i += 1
                    try:
                        with winreg.OpenKey(key, ver) as sub:
                            for name in (
                                "INSTALLLOCATION",
                                "OFFICEINSTALLLOCATION",
                                "BASISINSTALLLOCATION",
                            ):
                                try:
                                    loc, _ = winreg.QueryValueEx(sub, name)
                                    if loc:
                                        exe = _soffice_from_dir(Path(loc))
                                        if exe:
                                            return str(exe)
                                except OSError:
                                    continue
                    except OSError:
                        continue
        except OSError:
            continue
    return None


def _find_libreoffice() -> str | None:
    # 1) 用户配置 / 环境变量
    for env in ("SOFFICE_PATH", "LIBREOFFICE_HOME", "LIBREOFFICE_PATH"):
        val = os.environ.get(env, "").strip().strip('"')
        if not val:
            continue
        p = Path(val)
        if p.is_file() and p.name.lower().startswith("soffice"):
            return str(p)
        exe = _soffice_from_dir(p)
        if exe:
            return str(exe)

    cfg = Path(__file__).resolve().parent.parent / "config" / "libreoffice.txt"
    if cfg.is_file():
        for line in cfg.read_text(encoding="utf-8", errors="ignore").splitlines():
            line = line.strip().strip('"')
            if not line or line.startswith("#"):
                continue
            p = Path(line)
            if p.is_file():
                return str(p)
            exe = _soffice_from_dir(p)
            if exe:
                return str(exe)

    # 2) PATH
    for name in ("soffice", "soffice.exe", "libreoffice"):
        p = shutil.which(name)
        if p:
            return p

    # 3) Windows 注册表（非默认盘符安装，如 D:\LibreOffice）
    reg = _libreoffice_from_registry()
    if reg:
        return reg

    # 4) 常见安装目录
    if sys.platform == "win32":
        bases: list[Path] = []
        for key in ("ProgramFiles", "ProgramFiles(x86)", "LocalAppData"):
            v = os.environ.get(key)
            if v:
                bases.append(Path(v))
        bases.extend(
            [
                Path(r"C:\Program Files"),
                Path(r"C:\Program Files (x86)"),
                Path(os.path.expanduser(r"~\AppData\Local\Programs")),
            ]
        )
        seen: set[str] = set()
        for base in bases:
            if not base.is_dir():
                continue
            patterns = [
                base / "LibreOffice" / "program" / "soffice.exe",
                base / "Libre Office" / "program" / "soffice.exe",
            ]
            try:
                patterns.extend(base.glob("LibreOffice*/program/soffice.exe"))
            except OSError:
                pass
            for cand in patterns:
                s = str(cand.resolve()) if cand.is_file() else ""
                if s and s not in seen:
                    seen.add(s)
                    return s

    return None


def _libreoffice_to_pdf(src: Path, dst: Path) -> Path:
    soffice = _find_libreoffice()
    if not soffice:
        raise RuntimeError(
            "未找到 LibreOffice（soffice）。\n"
            "若已安装：在项目 config\\libreoffice.txt 写入安装目录，例如：\n"
            "  D:\\LibreOffice\n"
            "或设置环境变量 LIBREOFFICE_HOME=D:\\LibreOffice"
        )
    outdir = dst.parent
    outdir.mkdir(parents=True, exist_ok=True)
    print(f"[Office->PDF] LibreOffice 正在转换: {src.name}", flush=True)
    subprocess.run(
        [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(outdir), str(src)],
        check=True,
        capture_output=True,
        text=True,
    )
    print("[Office->PDF] LibreOffice 转换完成", flush=True)
    generated = outdir / f"{src.stem}.pdf"
    if not generated.is_file():
        raise RuntimeError("LibreOffice 转换后未找到 PDF")
    if generated.resolve() != dst.resolve():
        shutil.move(str(generated), str(dst))
    return dst


def _windows_office_to_pdf(src: Path, dst: Path, kind: str) -> Path:
    try:
        import comtypes.client  # type: ignore
    except ImportError as err:
        raise RuntimeError("Windows 下可安装: pip install comtypes") from err

    dst.parent.mkdir(parents=True, exist_ok=True)
    label = {"word": "Word", "excel": "Excel", "ppt": "PowerPoint"}.get(kind, kind)
    print(f"[Office->PDF] Microsoft {label} 正在转换: {src.name}", flush=True)
    if kind == "word":
        app = comtypes.client.CreateObject("Word.Application")
        app.Visible = False
        try:
            doc = app.Documents.Open(str(src.resolve()))
            doc.SaveAs(str(dst.resolve()), FileFormat=17)
            doc.Close()
        finally:
            app.Quit()
    elif kind == "excel":
        app = comtypes.client.CreateObject("Excel.Application")
        app.Visible = False
        try:
            wb = app.Workbooks.Open(str(src.resolve()))
            wb.ExportAsFixedFormat(0, str(dst.resolve()))
            wb.Close(False)
        finally:
            app.Quit()
    elif kind == "ppt":
        app = comtypes.client.CreateObject("PowerPoint.Application")
        try:
            pres = app.Presentations.Open(str(src.resolve()), WithWindow=False)
            pres.SaveAs(str(dst.resolve()), 32)
            pres.Close()
        finally:
            app.Quit()
    else:
        raise ValueError(kind)
    print(f"[Office->PDF] Microsoft {label} 转换完成", flush=True)
    return dst


def word_to_pdf_simple(src: Path, dst: Path) -> Path:
    """
    不依赖 Word/LibreOffice：仅把 Word 正文写入 PDF（纯文本，无原版式/图片/表格样式）。
    仅支持 .docx 等 python-docx 可读格式，不支持旧版 .doc。
    """
    from docx import Document

    if src.suffix.lower() in (".doc",):
        raise RuntimeError("旧版 .doc 无法用简易模式，请安装 LibreOffice 或 Word，或另存为 .docx")

    docx = Document(src)
    pdf = fitz.open()
    page = pdf.new_page(width=595, height=842)
    margin_x, margin_y = 50, 50
    y = margin_y
    line_h = 14
    max_y = 842 - margin_y
    fontsize = 11

    def new_page() -> fitz.Page:
        return pdf.new_page(width=595, height=842)

    page = pdf[0]
    paras = [p for p in docx.paragraphs if p.text.strip()]
    for para in iter_progress(paras or docx.paragraphs, desc="Word->PDF", unit="段"):
        text = para.text.replace("\n", " ").strip()
        if not text:
            y += line_h // 2
            continue
        while text:
            if y > max_y:
                page = new_page()
                y = margin_y
            chunk = text[:90]
            if len(text) > 90:
                cut = chunk.rfind(" ")
                if cut > 20:
                    chunk = text[:cut]
            page.insert_text((margin_x, y), chunk, fontsize=fontsize)
            text = text[len(chunk) :].lstrip()
            y += line_h

    dst.parent.mkdir(parents=True, exist_ok=True)
    pdf.save(dst)
    pdf.close()
    return dst


def excel_to_pdf_simple(src: Path, dst: Path) -> Path:
    """简易 Excel→PDF：表格数据按行列写入 PDF，无原版式/图表。"""
    import openpyxl

    if src.suffix.lower() in (".xls",):
        raise RuntimeError("旧版 .xls 无法用简易模式，请安装 LibreOffice 或另存为 .xlsx")

    wb = openpyxl.load_workbook(src, data_only=True, read_only=True)
    pdf = fitz.open()
    page = pdf.new_page(width=842, height=595)  # 横向便于多列
    margin_x, margin_y = 40, 40
    y = margin_y
    line_h = 13
    col_w = 72
    fontsize = 9
    max_y = 555

    def new_page() -> fitz.Page:
        return pdf.new_page(width=842, height=595)

    try:
        sheets = list(wb.worksheets)
        for sheet in iter_progress(sheets, desc="Excel->PDF", unit="表"):
            title = f"[工作表] {sheet.title}"
            if y > max_y - 30:
                page = new_page()
                y = margin_y
            page.insert_text((margin_x, y), title, fontsize=12)
            y += 20

            for row in sheet.iter_rows(values_only=True):
                if not any(c is not None and str(c).strip() for c in row):
                    continue
                if y > max_y:
                    page = new_page()
                    y = margin_y
                x = margin_x
                for cell in row:
                    val = "" if cell is None else str(cell).replace("\n", " ")
                    if len(val) > 14:
                        val = val[:14] + "…"
                    page.insert_text((x, y), val, fontsize=fontsize)
                    x += col_w
                    if x > 780:
                        break
                y += line_h
            y += 18
    finally:
        wb.close()

    dst.parent.mkdir(parents=True, exist_ok=True)
    pdf.save(dst)
    pdf.close()
    return dst


def ppt_to_pdf_simple(src: Path, dst: Path) -> Path:
    """简易 PPT→PDF：每页幻灯片提取文字与嵌入图片，版式不保证。"""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    if src.suffix.lower() in (".ppt",):
        raise RuntimeError("旧版 .ppt 无法用简易模式，请安装 LibreOffice 或另存为 .pptx")

    prs = Presentation(src)
    pdf = fitz.open()
    page_w, page_h = 595, 842

    slides = list(prs.slides)
    for idx, slide in enumerate(
        iter_progress(slides, desc="PPT->PDF", unit="页"), start=1
    ):
        page = pdf.new_page(width=page_w, height=page_h)
        y = 45
        page.insert_text((40, 20), f"幻灯片 {idx}", fontsize=13)

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img = shape.image
                    rect = fitz.Rect(40, y, 280, y + 160)
                    page.insert_image(rect, stream=img.blob)
                    y += 175
                except Exception:
                    pass
                continue

            text = getattr(shape, "text", "") or ""
            text = text.strip()
            if not text:
                continue
            for line in text.splitlines():
                line = line.strip()
                if not line:
                    continue
                if y > page_h - 50:
                    break
                chunk = line[:95]
                page.insert_text((40, y), chunk, fontsize=11)
                y += 15

    if len(prs.slides) == 0:
        pdf.new_page(width=page_w, height=page_h)

    dst.parent.mkdir(parents=True, exist_ok=True)
    pdf.save(dst)
    pdf.close()
    return dst


def office_to_pdf(src: Path, dst: Path | None = None) -> Path:
    ext = src.suffix.lower()
    if ext not in OFFICE_EXT:
        raise ValueError(f"不支持的 Office 格式: {ext}")
    kind = OFFICE_EXT[ext]
    dst = dst or _default_out(src, ".pdf")

    print(f"[Office->PDF] 开始: {src.name}", flush=True)

    try:
        return _libreoffice_to_pdf(src, dst)
    except Exception as exc:
        print(f"[Office->PDF] LibreOffice 不可用或未成功: {exc}", flush=True)

    if sys.platform == "win32":
        try:
            return _windows_office_to_pdf(src, dst, kind)
        except Exception as exc:
            print(f"[Office->PDF] Microsoft Office 未成功: {exc}", flush=True)

    simple = {
        "word": word_to_pdf_simple,
        "excel": excel_to_pdf_simple,
        "ppt": ppt_to_pdf_simple,
    }
    if kind in simple:
        print("[Office->PDF] 使用简易模式（纯 Python）...", flush=True)
        out = simple[kind](src, dst)
        print("[Office->PDF] 简易模式完成", flush=True)
        return out

    raise RuntimeError("无法转为 PDF，请检查文件格式。")


def file_to_pdf(src: Path, dst: Path | None = None) -> Path:
    ext = src.suffix.lower()
    dst = dst or _default_out(src, ".pdf")
    if ext in IMAGE_EXT:
        return images_to_pdf([src], dst)
    if ext in OFFICE_EXT:
        return office_to_pdf(src, dst)
    raise ValueError(f"无法转为 PDF: {ext}")


def encrypt_pdf(
    src: Path,
    dst: Path | None = None,
    user_password: str = "",
    owner_password: str | None = None,
) -> Path:
    dst = dst or src.with_name(f"{src.stem}_encrypted.pdf")
    doc = fitz.open(src)
    try:
        owner = owner_password if owner_password is not None else user_password
        doc.save(
            dst,
            encryption=fitz.PDF_ENCRYPT_AES_256,
            user_pw=user_password or "",
            owner_pw=owner or "",
        )
    finally:
        doc.close()
    return dst


def decrypt_pdf(
    src: Path,
    dst: Path | None = None,
    password: str = "",
) -> Path:
    dst = dst or src.with_name(f"{src.stem}_decrypted.pdf")
    doc = _open_pdf(src, password)
    try:
        doc.save(dst, garbage=4, deflate=True)
    finally:
        doc.close()
    return dst
